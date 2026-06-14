"""Deterministic post-process safety net for the deliverables.

The agent prompts forbid em/en dashes in client copy, but a model can still slip
one in. This pass programmatically replaces typographic dashes with a plain
hyphen-minus ("-") in the rendered report.html and report.pptx, so the shipped
files never contain one regardless of what the LLM did.
"""

from __future__ import annotations

from pathlib import Path

# All "fancy" dashes a model might emit -> hyphen-minus. Hyphen-minus ("-",
# U+002D) is intentionally left alone (ranges like "18-24" stay correct).
_DASHES = {
    "—": "-",  # — em dash
    "–": "-",  # – en dash
    "―": "-",  # ― horizontal bar
    "‒": "-",  # ‒ figure dash
}
_TRANS = str.maketrans(_DASHES)


def sanitize_text(s: str) -> tuple[str, int]:
    """Return (cleaned, count_replaced)."""
    if not s:
        return s, 0
    count = sum(s.count(ch) for ch in _DASHES)
    if not count:
        return s, 0
    return s.translate(_TRANS), count


def sanitize_html(path: str | Path) -> int:
    """Replace fancy dashes in an HTML file in place. Returns count replaced.

    A whole-file replace is safe: the dash code points never appear in base64
    data URIs (their alphabet is [A-Za-z0-9+/=]) or in HTML/CSS syntax.
    """
    path = Path(path)
    text = path.read_text(encoding="utf-8")
    cleaned, count = sanitize_text(text)
    if count:
        path.write_text(cleaned, encoding="utf-8")
    return count


def sanitize_pptx(path: str | Path) -> int:
    """Replace fancy dashes in every text run of a .pptx in place (run-level so
    formatting is preserved). Returns count replaced."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    path = Path(path)
    prs = Presentation(str(path))
    total = 0

    def fix_text_frame(tf) -> None:
        nonlocal total
        for para in tf.paragraphs:
            for run in para.runs:
                cleaned, count = sanitize_text(run.text)
                if count:
                    run.text = cleaned
                    total += count

    def walk(shapes) -> None:
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes)
                continue
            if shape.has_text_frame:
                fix_text_frame(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        fix_text_frame(cell.text_frame)

    for slide in prs.slides:
        walk(slide.shapes)

    if total:
        prs.save(str(path))
    return total


def sanitize_artifacts(html_path: str | Path | None = None,
                       pptx_path: str | Path | None = None) -> dict[str, int]:
    """Sanitize the report deliverables. Each format is best-effort and isolated
    so one failing file never blocks the other. Returns per-format counts."""
    out = {"html": 0, "pptx": 0}
    if html_path and Path(html_path).exists():
        out["html"] = sanitize_html(html_path)
    if pptx_path and Path(pptx_path).exists():
        out["pptx"] = sanitize_pptx(pptx_path)
    return out
